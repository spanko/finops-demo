"""
PowerPoint findings deck generator.

Produces the 10-12 slide findings deck described in Section 8 of the playbook.
Built from scratch using python-pptx (no external template required) in GT
brand colors.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

import duckdb
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from gt_finops.aggregate import summarize, top_findings


# GT brand palette
PURPLE_DARK = RGBColor(0x3A, 0x1D, 0x5C)
PURPLE      = RGBColor(0x58, 0x2C, 0x83)
PURPLE_LT   = RGBColor(0x95, 0x81, 0xB2)
PURPLE_SOFT = RGBColor(0xE8, 0xE2, 0xF0)
ORANGE      = RGBColor(0xF2, 0x6B, 0x23)
TEAL        = RGBColor(0x00, 0x83, 0x8F)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK   = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_BODY   = RGBColor(0x33, 0x33, 0x33)
TEXT_MUTED  = RGBColor(0x76, 0x76, 0x76)
LINE        = RGBColor(0xD1, 0xD1, 0xD1)
CARD_BG     = RGBColor(0xF7, 0xF5, 0xFA)

# Category colors — use the GT two-view palette consistently
# M365 + Security = cost view (orange family); Commitments + Waste = architecture/operational (teal family)
# But per the original deck design, we keep distinct hues for each category
CATEGORY_COLORS = {
    "M365":        RGBColor(0x58, 0x2C, 0x83),  # GT primary purple
    "Security":    RGBColor(0x95, 0x81, 0xB2),  # GT lavender
    "Commitments": RGBColor(0x00, 0x83, 0x8F),  # GT teal
    "Waste":       RGBColor(0xF2, 0x6B, 0x23),  # GT orange
}


def write_findings_deck(
    conn: duckdb.DuckDBPyConnection,
    output_path: Path,
    client_name: str,
    engagement_start: str | None = None,
    engagement_end: str | None = None,
) -> Path:
    """Build the findings deck as a .pptx file."""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # 6 = blank layout

    summary = summarize(conn)
    top = top_findings(conn, limit=10)

    dates_str = ""
    if engagement_start and engagement_end:
        dates_str = f"{engagement_start} — {engagement_end}"
    elif engagement_end:
        dates_str = f"through {engagement_end}"
    else:
        dates_str = datetime.utcnow().strftime("%Y-%m-%d")

    # Slides
    _add_cover_slide(prs, blank, client_name, dates_str)
    _add_summary_slide(prs, blank, client_name, summary)
    _add_headline_slide(prs, blank, summary)
    _add_category_slide(prs, blank, summary, "M365", "M365 Licensing findings", ORANGE)
    _add_category_slide(prs, blank, summary, "Security", "Microsoft Security findings", TEAL)
    _add_category_slide(prs, blank, summary, "Commitments", "Azure commitment findings", PURPLE)
    _add_category_slide(prs, blank, summary, "Waste", "Azure waste findings", PURPLE_LT)
    _add_top_findings_slide(prs, blank, top)
    _add_roadmap_slide(prs, blank, summary)
    _add_commercial_slide(prs, blank)
    _add_continuous_slide(prs, blank)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


# -----------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------

def _bg(slide, color: RGBColor) -> None:
    """Set full slide background color."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def _rect(slide, x, y, w, h, fill: RGBColor, line_color: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _text_box(slide, x, y, w, h, text, *, font="Arial", size=14,
              color=TEXT_BODY, bold=False, italic=False,
              align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return tb


def _kicker(slide, x, y, text, color=PURPLE):
    _text_box(slide, x, y, Inches(8), Inches(0.3),
              text.upper(), size=10, bold=True, color=color)


def _title_bar(slide):
    """Thin purple bar at the top edge — subtle brand element."""
    _rect(slide, 0, 0, Inches(13.333), Inches(0.1), PURPLE)


def _footer(slide, page_num):
    _text_box(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.25),
              "Grant Thornton · FinOps Findings",
              size=8, color=TEXT_MUTED)
    _text_box(slide, Inches(12.5), Inches(7.1), Inches(0.5), Inches(0.25),
              str(page_num), size=8, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


# -----------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------

def _add_cover_slide(prs, layout, client_name, dates_str):
    slide = prs.slides.add_slide(layout)
    _bg(slide, PURPLE_DARK)

    # Orange + teal accent squares (two-view motif)
    _rect(slide, Inches(0.6), Inches(0.6), Inches(0.25), Inches(0.25), ORANGE)
    _rect(slide, Inches(0.92), Inches(0.6), Inches(0.25), Inches(0.25), TEAL)

    _text_box(slide, Inches(0.6), Inches(1.1), Inches(10), Inches(0.4),
              "GRANT THORNTON  ·  FINOPS ENGAGEMENT FINDINGS",
              size=12, bold=True, color=RGBColor(0xC5, 0xB8, 0xD6))

    _text_box(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(1.3),
              client_name, font="Arial", size=60, bold=True, color=WHITE)

    _text_box(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.6),
              "Cost optimization analysis", size=26, color=RGBColor(0xCB, 0xC0, 0xD8))

    _text_box(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
              dates_str, size=16, italic=True, color=ORANGE)

    _text_box(slide, Inches(0.6), Inches(6.9), Inches(12), Inches(0.3),
              "Engagement findings deck · Internal Grant Thornton",
              size=10, color=RGBColor(0xA0, 0x93, 0xB5))


def _add_summary_slide(prs, layout, client_name, summary):
    slide = prs.slides.add_slide(layout)
    _title_bar(slide)
    _kicker(slide, Inches(0.6), Inches(0.5), "Engagement summary")
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7),
              f"We analyzed {client_name}'s Microsoft ecosystem",
              size=30, bold=True, color=PURPLE_DARK)

    stats = [
        (f"{summary.total_findings}", "findings surfaced"),
        (f"{len(summary.by_recipe)}", "recipes executed"),
        (f"${summary.total_capturable_usd/1000:,.0f}K", "capturable annual savings"),
    ]

    x = Inches(0.6)
    y = Inches(2.5)
    card_w = Inches(3.9)
    card_h = Inches(2.2)
    gap = Inches(0.25)

    for i, (value, label) in enumerate(stats):
        cx = x + (card_w + gap) * i
        _rect(slide, cx, y, card_w, card_h, CARD_BG)
        accent_color = [PURPLE, ORANGE, TEAL][i]
        _rect(slide, cx, y, Inches(0.08), card_h, accent_color)
        _text_box(slide, cx + Inches(0.35), y + Inches(0.3),
                  card_w - Inches(0.5), Inches(0.3),
                  label.upper(), size=10, bold=True, color=accent_color)
        _text_box(slide, cx + Inches(0.35), y + Inches(0.7),
                  card_w - Inches(0.5), Inches(1.2),
                  value, size=60, bold=True, color=TEXT_DARK)

    _text_box(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.9),
              "Findings are categorized by M365 Licensing, Microsoft Security, Azure Commitments, and Azure Waste.\n"
              "Each finding carries a gross savings estimate, a capturable estimate, and days-to-capture. "
              "The numbers reported throughout this deck are the capturable estimates.",
              size=14, color=TEXT_BODY)

    _footer(slide, 2)


def _add_headline_slide(prs, layout, summary):
    slide = prs.slides.add_slide(layout)
    _title_bar(slide)
    _kicker(slide, Inches(0.6), Inches(0.5), "Headline numbers")
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7),
              "Savings by category",
              size=30, bold=True, color=PURPLE_DARK)

    # Build a horizontal bar-style layout per category
    y = Inches(2.0)
    row_h = Inches(0.85)

    for i, (cat, stats) in enumerate(summary.by_category.items()):
        ry = y + row_h * i
        cat_color = CATEGORY_COLORS.get(cat, PURPLE)

        # Color badge
        _rect(slide, Inches(0.6), ry, Inches(1.6), Inches(0.55), cat_color)
        _text_box(slide, Inches(0.6), ry, Inches(1.6), Inches(0.55),
                  cat.upper(), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Count
        _text_box(slide, Inches(2.4), ry + Inches(0.07),
                  Inches(1.8), Inches(0.45),
                  f"{stats['count']} findings", size=14, color=TEXT_BODY)

        # Savings
        _text_box(slide, Inches(4.4), ry, Inches(3), Inches(0.3),
                  "CAPTURABLE ANNUAL", size=9, bold=True, color=TEXT_MUTED)
        _text_box(slide, Inches(4.4), ry + Inches(0.2),
                  Inches(4), Inches(0.5),
                  f"${stats['capturable']:,.0f}",
                  size=26, bold=True, color=TEXT_DARK)

        _text_box(slide, Inches(8.8), ry, Inches(3), Inches(0.3),
                  "GROSS IF FULLY CAPTURED", size=9, bold=True, color=TEXT_MUTED)
        _text_box(slide, Inches(8.8), ry + Inches(0.2),
                  Inches(4), Inches(0.5),
                  f"${stats['gross']:,.0f}",
                  size=18, color=TEXT_MUTED)

    # Total row
    y_total = y + row_h * len(summary.by_category) + Inches(0.3)
    _rect(slide, Inches(0.6), y_total, Inches(12), Inches(0.8), PURPLE_DARK)
    _text_box(slide, Inches(0.9), y_total + Inches(0.15),
              Inches(4), Inches(0.5),
              "TOTAL CAPTURABLE", size=11, bold=True, color=ORANGE)
    _text_box(slide, Inches(0.9), y_total + Inches(0.35),
              Inches(6), Inches(0.5),
              f"${summary.total_capturable_usd:,.0f}/year",
              size=22, bold=True, color=WHITE)

    _footer(slide, 3)


def _add_category_slide(prs, layout, summary, category, title, accent_color):
    slide = prs.slides.add_slide(layout)
    _title_bar(slide)
    _kicker(slide, Inches(0.6), Inches(0.5), f"Category · {category}", color=accent_color)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7),
              title, size=30, bold=True, color=PURPLE_DARK)

    # List recipes in this category with their counts and totals
    recipes_in_cat = []
    for rid, stats in summary.by_recipe.items():
        # Look up category by querying — we don't have it here directly.
        # Easy workaround: check the by_category totals column match.
        pass

    # Easier: iterate recipes and filter by category via find-in-dict
    matching = []
    # Need to query: but we only have summary. Use the stored recipe registry.
    from gt_finops.recipes import RECIPE_BY_ID
    for rid, stats in summary.by_recipe.items():
        cls = RECIPE_BY_ID.get(rid)
        if cls and cls.category == category:
            matching.append((rid, stats))

    if not matching:
        _text_box(slide, Inches(0.6), Inches(2.5), Inches(10), Inches(0.6),
                  "No findings in this category.",
                  size=16, italic=True, color=TEXT_MUTED)
        _footer(slide, 4)
        return

    y = Inches(2.0)
    row_h = Inches(0.7)
    for i, (rid, stats) in enumerate(matching):
        ry = y + row_h * i
        if ry > Inches(6.5):
            break  # don't overflow
        _rect(slide, Inches(0.6), ry, Inches(0.1), Inches(0.5), accent_color)
        _text_box(slide, Inches(0.85), ry + Inches(0.05),
                  Inches(1.0), Inches(0.4),
                  rid, size=12, bold=True, color=accent_color)
        _text_box(slide, Inches(1.7), ry + Inches(0.05),
                  Inches(6.5), Inches(0.4),
                  stats["name"], size=13, bold=True, color=TEXT_DARK)
        _text_box(slide, Inches(8.5), ry + Inches(0.05),
                  Inches(2), Inches(0.4),
                  f"{stats['count']} findings", size=11, color=TEXT_MUTED)
        _text_box(slide, Inches(10.5), ry + Inches(0.05),
                  Inches(2.5), Inches(0.4),
                  f"${stats['capturable']:,.0f}",
                  size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.RIGHT)

    # Total footer
    cat_total = summary.by_category.get(category, {}).get("capturable", 0)
    _rect(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5), CARD_BG)
    _text_box(slide, Inches(0.9), Inches(6.55), Inches(6), Inches(0.4),
              f"{category} subtotal capturable", size=12, bold=True, color=PURPLE_DARK)
    _text_box(slide, Inches(9), Inches(6.55), Inches(3.5), Inches(0.4),
              f"${cat_total:,.0f}/year", size=14, bold=True,
              color=accent_color, align=PP_ALIGN.RIGHT)

    _footer(slide, 4)


def _add_top_findings_slide(prs, layout, top_df):
    slide = prs.slides.add_slide(layout)
    _title_bar(slide)
    _kicker(slide, Inches(0.6), Inches(0.5), "Biggest individual wins")
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7),
              "Top 10 findings by capturable savings",
              size=30, bold=True, color=PURPLE_DARK)

    # Table header
    y = Inches(2.0)
    _rect(slide, Inches(0.6), y, Inches(12), Inches(0.4), PURPLE)
    headers = [("#", Inches(0.6), Inches(0.5)),
               ("Recipe", Inches(1.1), Inches(1.0)),
               ("Entity", Inches(2.2), Inches(3.5)),
               ("Recommendation", Inches(5.8), Inches(5.0)),
               ("Capturable", Inches(10.9), Inches(1.7))]
    for label, x, w in headers:
        _text_box(slide, x + Inches(0.1), y + Inches(0.08), w, Inches(0.3),
                  label, size=11, bold=True, color=WHITE)

    row_h = Inches(0.42)
    for i, row in enumerate(top_df.to_dict("records")):
        ry = y + Inches(0.4) + row_h * i
        if i % 2 == 0:
            _rect(slide, Inches(0.6), ry, Inches(12), row_h, CARD_BG)
        _text_box(slide, Inches(0.7), ry + Inches(0.08), Inches(0.4), Inches(0.3),
                  str(i + 1), size=11, color=TEXT_DARK)
        _text_box(slide, Inches(1.2), ry + Inches(0.08), Inches(0.9), Inches(0.3),
                  row["recipe_id"], size=10, color=TEXT_BODY)
        _text_box(slide, Inches(2.3), ry + Inches(0.08), Inches(3.4), Inches(0.3),
                  row["entity_name"][:42], size=10, color=TEXT_DARK)
        _text_box(slide, Inches(5.9), ry + Inches(0.08), Inches(4.9), Inches(0.3),
                  row["recommended_state"][:60], size=10, color=TEXT_BODY)
        _text_box(slide, Inches(11.0), ry + Inches(0.08), Inches(1.5), Inches(0.3),
                  f"${row['capturable_annual_savings_usd']:,.0f}",
                  size=11, bold=True, color=TEXT_DARK, align=PP_ALIGN.RIGHT)

    _footer(slide, 8)


def _add_roadmap_slide(prs, layout, summary):
    slide = prs.slides.add_slide(layout)
    _title_bar(slide)
    _kicker(slide, Inches(0.6), Inches(0.5), "Implementation roadmap")
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7),
              "How we sequence the work",
              size=30, bold=True, color=PURPLE_DARK)

    phases = [
        ("WEEK 1", "Low-risk, high-ROI wins",
         "Disabled account license cleanup · Orphan resources · AHB audit entry points",
         ORANGE),
        ("WEEKS 2-3", "License right-sizing",
         "E5→E3 review with BU leads · F3 candidate confirmation · Copilot reclaim and reassignment",
         PURPLE),
        ("WEEKS 3-4", "Security and commitments",
         "Defender P2→P1 review with security · Sentinel tier adjustment · RI re-scoping",
         TEAL),
        ("WEEK 4", "Commercial conversation",
         "Microsoft account team: findings-backed renewal discussion and SA reconciliation",
         PURPLE_LT),
    ]

    y = Inches(2.0)
    row_h = Inches(1.2)
    for i, (when, what, detail, color) in enumerate(phases):
        ry = y + row_h * i
        _rect(slide, Inches(0.6), ry, Inches(0.1), Inches(1.0), color)
        _text_box(slide, Inches(0.9), ry + Inches(0.1), Inches(1.6), Inches(0.3),
                  when, size=11, bold=True, color=color, font="Arial")
        _text_box(slide, Inches(2.6), ry + Inches(0.05), Inches(10), Inches(0.4),
                  what, size=14, bold=True, color=TEXT_DARK)
        _text_box(slide, Inches(2.6), ry + Inches(0.4), Inches(10), Inches(0.6),
                  detail, size=11, color=TEXT_BODY)

    _footer(slide, 9)


def _add_commercial_slide(prs, layout):
    slide = prs.slides.add_slide(layout)
    _title_bar(slide)
    _kicker(slide, Inches(0.6), Inches(0.5), "Commercial recommendations")
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7),
              "Use the findings to improve your Microsoft position",
              size=30, bold=True, color=PURPLE_DARK)

    items = [
        ("EA renewal leverage",
         "A credible reclaim list creates negotiating position ahead of EA renewal. "
         "Microsoft reps often meet halfway with better pricing rather than accept seat count drop."),
        ("Copilot commitment repricing",
         "Over-committed Copilot licenses can be renegotiated using actual adoption data."),
        ("Azure consumption commitment accuracy",
         "Accurate 12-month forecasts replace vague growth estimates — better tiering unlocked."),
        ("SA entitlement validation",
         "AHB findings often surface SA entitlement mismatches. Reconciliation with rep "
         "improves licensing posture without new spend."),
    ]

    y = Inches(2.0)
    row_h = Inches(1.15)
    for i, (title, body) in enumerate(items):
        ry = y + row_h * i
        _rect(slide, Inches(0.6), ry, Inches(12), Inches(1.0), CARD_BG)
        _rect(slide, Inches(0.6), ry, Inches(0.08), Inches(1.0), ORANGE)
        _text_box(slide, Inches(0.85), ry + Inches(0.12), Inches(11.5), Inches(0.35),
                  title, size=14, bold=True, color=PURPLE_DARK)
        _text_box(slide, Inches(0.85), ry + Inches(0.45), Inches(11.5), Inches(0.55),
                  body, size=11, color=TEXT_BODY)

    _footer(slide, 10)


def _add_continuous_slide(prs, layout):
    slide = prs.slides.add_slide(layout)
    _bg(slide, PURPLE_DARK)

    _rect(slide, Inches(0.6), Inches(0.6), Inches(0.25), Inches(0.25), ORANGE)
    _rect(slide, Inches(0.92), Inches(0.6), Inches(0.25), Inches(0.25), TEAL)

    _text_box(slide, Inches(0.6), Inches(1.1), Inches(10), Inches(0.4),
              "BEYOND THIS ENGAGEMENT",
              size=12, bold=True, color=RGBColor(0xC5, 0xB8, 0xD6))

    _text_box(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(1.0),
              "From a one-time exercise",
              size=36, bold=True, color=WHITE)
    _text_box(slide, Inches(0.6), Inches(2.6), Inches(12), Inches(1.0),
              "to continuous optimization",
              size=36, bold=True, color=ORANGE)

    _text_box(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(0.4),
              "Quarterly re-runs of this same analysis", size=16, bold=True,
              color=WHITE)
    _text_box(slide, Inches(0.6), Inches(4.4), Inches(12), Inches(0.5),
              "Findings decay — re-ingest and re-analyze every 90 days.",
              size=13, color=RGBColor(0xC5, 0xB8, 0xD6))

    _text_box(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
              "Agentic AI — Ask, Surface, Act", size=16, bold=True, color=WHITE)
    _text_box(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
              "Always-on platform engagement: conversational Q&A against FinOps warehouse, "
              "scheduled sweeps, gated remediation.",
              size=13, color=RGBColor(0xC5, 0xB8, 0xD6))

    _text_box(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4),
              "Grant Thornton continues as the optimization partner.",
              size=14, italic=True, color=ORANGE)
